"""
utils.py - Knowledge base document parsing.

Reads every file in the knowledge_base/ folder and extracts plain text.
Supported: PDF, DOCX, TXT/MD/CSV, plus (new) XLSX, PPTX, HTML.

The original app only handled PDF/DOCX/TXT. The extra parsers are additive
and degrade gracefully if the optional library isn't installed.
"""

import os

import fitz  # PyMuPDF
import docx


def parse_file(file_path: str) -> str:
    """Extracts text from a given file based on its extension."""
    if not os.path.exists(file_path):
        return ""

    ext = file_path.lower().split(".")[-1]

    if ext == "pdf":
        return parse_pdf(file_path)
    elif ext in ["doc", "docx"]:
        return parse_docx(file_path)
    elif ext in ["txt", "md", "csv"]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif ext in ["xlsx", "xlsm"]:
        return parse_xlsx(file_path)
    elif ext in ["pptx"]:
        return parse_pptx(file_path)
    elif ext in ["html", "htm"]:
        return parse_html(file_path)
    else:
        return f"[Media File: {os.path.basename(file_path)} cannot be parsed locally yet]"


def parse_pdf(file_path: str) -> str:
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
    except Exception as e:
        print(f"Error parsing PDF {file_path}: {e}")
    return text


def parse_docx(file_path: str) -> str:
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Also pull text out of tables, which docx.paragraphs skips.
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text += " | ".join(cells) + "\n"
    except Exception as e:
        print(f"Error parsing DOCX {file_path}: {e}")
    return text


def parse_xlsx(file_path: str) -> str:
    """Extract text from a spreadsheet, sheet by sheet."""
    text = ""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            text += f"\n[Sheet: {sheet.title}]\n"
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text += " | ".join(cells) + "\n"
        wb.close()
    except Exception as e:
        print(f"Error parsing XLSX {file_path}: {e}")
    return text


def parse_pptx(file_path: str) -> str:
    """Extract text from every shape on every slide."""
    text = ""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            text += f"\n[Slide {i}]\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            text += line + "\n"
    except Exception as e:
        print(f"Error parsing PPTX {file_path}: {e}")
    return text


def parse_html(file_path: str) -> str:
    """Strip tags and return readable text."""
    text = ""
    try:
        from bs4 import BeautifulSoup

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
    except Exception as e:
        print(f"Error parsing HTML {file_path}: {e}")
    return text


def get_all_kb_text(kb_dir: str) -> str:
    """Reads all files in the KB directory and combines their text."""
    if not os.path.exists(kb_dir):
        return ""

    combined_text = ""
    for filename in sorted(os.listdir(kb_dir)):
        filepath = os.path.join(kb_dir, filename)
        if os.path.isfile(filepath) and not filename.startswith("."):
            combined_text += f"\n\n--- DOCUMENT: {filename} ---\n"
            combined_text += parse_file(filepath)

    return combined_text
